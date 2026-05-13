from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Function to add a slide with title and content
    def add_slide(title_text, content_text):
        slide_layout = prs.slide_layouts[1] # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = title_text
        content.text = content_text

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "VOTE HUB"
    subtitle.text = "Secure & Transparent Online Voting System\nPresented By: Ayush\nProject: Web-Based Voting Portal"

    # Slide 2: Introduction
    add_slide("Introduction", 
              "• Vote Hub is a web-based digital voting platform.\n"
              "• Aims to digitize the manual voting process.\n"
              "• Focuses on accessibility, speed, and transparency.")

    # Slide 3: Problem Statement
    add_slide("Problem Statement",
              "• Manual voting is time-consuming.\n"
              "• High risk of human error in counting.\n"
              "• Logistics and printing costs are high.\n"
              "• Low turnout due to physical presence requirements.")

    # Slide 4: Proposed Solution
    add_slide("Proposed Solution",
              "• Online portal for instant voting.\n"
              "• Automated counting for real-time results.\n"
              "• Paperless and eco-friendly system.\n"
              "• Accessible from anywhere via the internet.")

    # Slide 5: Objectives
    add_slide("Objectives",
              "• Ensure 'One Voter, One Vote' policy.\n"
              "• Maintain 100% data accuracy.\n"
              "• Provide a user-friendly and intuitive UI.\n"
              "• Securely store voter and party information.")

    # Slide 6: System Architecture
    add_slide("System Architecture",
              "• Client-Server Model.\n"
              "• Frontend: HTML/CSS/JS for user interaction.\n"
              "• Backend: PHP for processing logic.\n"
              "• Database: MySQL for secure data storage.")

    # Slide 7: Tech Stack
    add_slide("Technology Stack",
              "• Languages: PHP, HTML5, CSS3, JavaScript.\n"
              "• Icons: Lucide Icons.\n"
              "• Database: MySQL.\n"
              "• Environment: XAMPP / WAMP Server.")

    # Slide 8: Database Design
    add_slide("Database Design",
              "• Users Table: Stores credentials and 'has_voted' flag.\n"
              "• Parties Table: Stores party names and vote counts.\n"
              "• Unique Constraints: Email-based registration to prevent duplicates.")

    # Slide 9: Module - Voter Portal
    add_slide("Module: Voter Portal",
              "• Secure Registration & Login.\n"
              "• Interactive Voting Dashboard.\n"
              "• Instant feedback after casting a vote.\n"
              "• Post-voting lockout to prevent double voting.")

    # Slide 10: Module - Admin Portal
    add_slide("Module: Admin Portal",
              "• Restricted access for system owners.\n"
              "• Real-time visualization of results.\n"
              "• System monitoring and data management.")

    # Slide 11: Security Features
    add_slide("Security Features",
              "• Bcrypt Password Hashing.\n"
              "• Session-based authentication.\n"
              "• Prevention of SQL Injection using PDO.")

    # Slide 12: Preventing Double Voting
    add_slide("Preventing Double Voting",
              "• Logic: System checks 'has_voted' status before displaying parties.\n"
              "• Database Transactions: Ensures atomicity during vote updates.\n"
              "• Session Validation: Prevents URL-based bypass.")

    # Slide 13: UI/UX Design
    add_slide("UI/UX Design Strategy",
              "• Modern Dark/Light theme support.\n"
              "• Mobile-responsive layout.\n"
              "• Clear Call-to-Action (CTA) buttons.")

    # Slide 14: Key Implementation (Code)
    add_slide("Implementation Highlights",
              "• Transaction handling for vote integrity.\n"
              "• Dynamic data fetching using PHP PDO.\n"
              "• Password encryption for user privacy.")

    # Slide 15: Testing Phase
    add_slide("Testing & Validation",
              "• Unit Testing: Individual form validations.\n"
              "• Integration Testing: End-to-end flow from login to vote.\n"
              "• Security Audit: Protecting against unauthorized access.")

    # Slide 16: Advantages
    add_slide("Key Advantages",
              "• Significant reduction in voting time.\n"
              "• Cost-effective (Zero paper/logistics).\n"
              "• Highly scalable for large populations.")

    # Slide 17: Challenges
    add_slide("Challenges Faced",
              "• Managing concurrent database updates.\n"
              "• Ensuring a seamless UI across different browsers.\n"
              "• Implementing secure session timeouts.")

    # Slide 18: Future Enhancements
    add_slide("Future Enhancements",
              "• Blockchain for immutable vote records.\n"
              "• Biometric (Fingerprint/Face) Auth.\n"
              "• SMS-based OTP verification.")

    # Slide 19: Conclusion
    add_slide("Conclusion",
              "• Vote Hub successfully modernizes the electoral process.\n"
              "• Provides a robust foundation for digital democracy.\n"
              "• Scalable, secure, and ready for deployment.")

    # Slide 20: Thank You
    add_slide("Thank You", "Questions?\n\nProject: Vote Hub\nStatus: Completed")

    # Save the presentation
    prs.save('Vote_Hub_Presentation.pptx')

if __name__ == "__main__":
    create_presentation()
